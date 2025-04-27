import PyPDF2
import os
from pptx import Presentation

def extract_text_from_pdf(pdf_path):
    with open(pdf_path, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        text = ""

        # Get the number of pages
        num_pages = len(pdf_reader.pages)
        print(f"Total pages: {num_pages}")

        # Extract text from each page
        for page_num in range(num_pages):
            page = pdf_reader.pages[page_num]
            text += f"\n\n--- Page {page_num + 1} ---\n\n"
            text += page.extract_text()

        return text

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = ""

    # Get the number of slides
    num_slides = len(prs.slides)
    print(f"Total slides: {num_slides}")

    # Extract text from each slide
    for slide_num, slide in enumerate(prs.slides):
        text += f"\n\n--- Slide {slide_num + 1} ---\n\n"

        # Extract text from shapes in the slide
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text += shape.text + "\n"

    return text

def process_file(file_path):
    file_name = os.path.basename(file_path)
    output_file = f"{os.path.splitext(file_name)[0]}_content.txt"

    print(f"Processing {file_path}...")

    if file_path.lower().endswith('.pdf'):
        extracted_text = extract_text_from_pdf(file_path)
    elif file_path.lower().endswith('.pptx'):
        extracted_text = extract_text_from_pptx(file_path)
    else:
        print(f"Unsupported file format: {file_path}")
        return

    # Save the extracted text to a file
    with open(output_file, "w", encoding="utf-8") as text_file:
        text_file.write(extracted_text)

    print(f"Text extraction complete. Saved to {output_file}")

if __name__ == "__main__":
    files_to_process = [
        "week2.pdf",
        "week3.pdf",
        "week4.pdf",
        "week5.pdf",
        "week6.pdf",
        "week7.pdf",
        "week8.pptx",
        "week9.pptx",
        "week10.pdf"
    ]

    for file_path in files_to_process:
        process_file(file_path)
